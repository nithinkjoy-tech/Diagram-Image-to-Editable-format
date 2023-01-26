from pptx import Presentation
from pptx.util import Pt,Inches
from pptx.dml.color import RGBColor
import sys
import cv2
import pytesseract
import subprocess
# Open the PowerPoint file
no=int(sys.argv[1])
fname=sys.argv[2]
prs = Presentation('./shapes-test.pptx')

# Get the first slide
slide = prs.slides[0]

# Get the first shape on the slide
for i in range(1,no+1):
    shape = slide.shapes[i]
    shape.line.width = Pt(2)
    shape.line.color.rgb = RGBColor(0, 0, 0)

# Save the changes to the PowerPoint file
prs.slide_width = Inches(12)
prs.slide_height = Inches(10)
prs.save('./new.pptx')



image = cv2.imread(fname)

# Convert the image to grayscale
gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

# Run OCR on the image
text = pytesseract.image_to_string(gray)
text=text.split("\n")
while("" in text):
    text.remove("")
# print(text)

    # print(len(ele))
# Get the location of the text in the image
location = pytesseract.image_to_boxes(gray)
xloc=[]
yloc=[]
char=[]
# Extract the x, y coordinates of the text
for box in location.splitlines():
    box = box.split(' ')
    x, y, w, h = int(box[1]), int(box[2]), int(box[3]), int(box[4])
    cv2.rectangle(image, (x,y), (x+w, y+h), (0, 0, 255), 2)
    char.append(box[0])
    xloc.append(x)
    yloc.append(y)
all=[char,xloc,yloc]
# subprocess.run(["node", "./addText.js", char,xloc,yloc])
print(all)

sys.stdout.flush()
