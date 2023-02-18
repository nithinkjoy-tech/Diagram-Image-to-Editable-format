try:
    import os
    os.environ['TF_CPP_MIN_LOG_LEVEL'] = '3' 
    import easyocr
    import cv2
    import numpy as np
    from pptx import Presentation
    from pptx.util import Pt,Inches
    from pptx.dml.color import RGBColor
    import sys
    import subprocess
    
    no=int(sys.argv[1])
    fname=sys.argv[2]
    prs = Presentation('./shapes-test.pptx')

    slide = prs.slides[0]

    for i in range(1,no+1):
        shape = slide.shapes[i]
        shape.line.width = Pt(2)
        shape.line.color.rgb = RGBColor(0, 0, 0)

    prs.slide_width = Inches(11.5)
    prs.slide_height = Inches(10)
    prs.save('./shapes-test.pptx')
except Exception as e:
    print(e)
    raise e
    sys.stdout.flush()