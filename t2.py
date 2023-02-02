import easyocr
import cv2
from matplotlib import pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Pt,Inches
from pptx.dml.color import RGBColor
import sys
import pytesseract
from pytesseract import Output
import subprocess
# Open the PowerPoint file
fname="bc.png"
no=6
prs = Presentation('./shapes-test.pptx')

# Get the first slide
slide = prs.slides[0]

# Get the first shape on the slide
for i in range(1,no+1):
    shape = slide.shapes[i]
    shape.line.width = Pt(2)
    shape.line.color.rgb = RGBColor(0, 0, 0)

# Save the changes to the PowerPoint file
prs.slide_width = Inches(11.5)
prs.slide_height = Inches(9.5)
prs.save('./new.pptx')

IMAGE_PATH = fname

reader = easyocr.Reader(['en'],gpu=True)
result = reader.readtext(IMAGE_PATH)

top_left = tuple(result[0][0][0])
bottom_right = tuple(result[0][0][2])
text = result[0][1]
font = cv2.FONT_HERSHEY_SIMPLEX

img = cv2.imread(IMAGE_PATH)
spacer = 100

alldata=[]
count=0
for detection in result: 
    top_left = tuple(detection[0][0])
    bottom_right = tuple(detection[0][2])
    text = detection[1]
    img = cv2.rectangle(img,top_left,bottom_right,(0,255,0),3)
    img = cv2.putText(img,text,(20,spacer), font, 0.5,(0,255,0),2,cv2.LINE_AA)
    spacer+=15
    loclist=list(top_left)
    count+=1
    alldata.append([text,loclist[0],loclist[1]])
    # xloc.append(loclist[0])
    # yloc.append(loclist[1])
print(alldata)
sys.stdout.flush()